import os
from pdf2image import convert_from_path
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from PyPDF2 import PdfReader
from tkinter import font as tkfont

class PDFConverterApp:
    def __init__(self, root):
        self.root = root
        root.title("PDF Converter")
        root.geometry("500x400")
        root.resizable(False, False)

        # ฟอนต์มน ๆ
        self.custom_font = tkfont.Font(family="Segoe UI", size=11)

        # Label header
        header = tk.Label(root, text="📚 PDF Converter", font=("Segoe UI", 16, "bold"))
        header.pack(pady=10)

        # PDF file
        self.pdf_path = tk.StringVar()
        pdf_frame = tk.Frame(root)
        pdf_frame.pack(pady=5, fill="x", padx=20)
        tk.Label(pdf_frame, text="PDF File:", font=self.custom_font).pack(side="left")
        tk.Entry(pdf_frame, textvariable=self.pdf_path, font=self.custom_font, width=35).pack(side="left", padx=5)
        tk.Button(pdf_frame, text="Browse", font=self.custom_font, command=self.browse_pdf).pack(side="left")

        # Page range
        self.page_range = tk.StringVar()
        page_frame = tk.Frame(root)
        page_frame.pack(pady=5, fill="x", padx=20)
        tk.Label(page_frame, text="Page Range (ex: 1-10):", font=self.custom_font).pack(side="left")
        tk.Entry(page_frame, textvariable=self.page_range, font=self.custom_font, width=15).pack(side="left", padx=5)

        # Export choice
        self.export_choice = tk.StringVar(value="1")
        export_frame = tk.Frame(root)
        export_frame.pack(pady=5, fill="x", padx=20)
        tk.Label(export_frame, text="Export Mode:", font=self.custom_font).pack(side="left")
        self.combo_export = ttk.Combobox(export_frame, values=["1 - Word", "2 - PowerPoint", "3 - Images Only", "1,2,3"], font=self.custom_font, width=20)
        self.combo_export.current(0)
        self.combo_export.pack(side="left", padx=5)

        # Output folder
        self.output_folder = tk.StringVar()
        folder_frame = tk.Frame(root)
        folder_frame.pack(pady=5, fill="x", padx=20)
        tk.Label(folder_frame, text="Output Folder:", font=self.custom_font).pack(side="left")
        tk.Entry(folder_frame, textvariable=self.output_folder, font=self.custom_font, width=30).pack(side="left", padx=5)
        tk.Button(folder_frame, text="Browse", font=self.custom_font, command=self.browse_folder).pack(side="left")

        # Convert button
        tk.Button(root, text="🚀 Convert", font=("Segoe UI", 13, "bold"), bg="#4CAF50", fg="white",
                  command=self.start_conversion).pack(pady=20)

        # Status
        self.status = tk.Label(root, text="", font=self.custom_font, fg="blue")
        self.status.pack(pady=5)

    def browse_pdf(self):
        file = filedialog.askopenfilename(title="Select PDF", filetypes=[("PDF files", "*.pdf")])
        if file:
            self.pdf_path.set(file)

    def browse_folder(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.output_folder.set(folder)

    def start_conversion(self):
        pdf_file = self.pdf_path.get()
        if not pdf_file or not os.path.isfile(pdf_file):
            messagebox.showerror("Error", "Please select a valid PDF file")
            return
        pdf_name = os.path.splitext(os.path.basename(pdf_file))[0]

        output_folder = self.output_folder.get()
        if not output_folder or not os.path.isdir(output_folder):
            messagebox.showerror("Error", "Please select a valid output folder")
            return

        # อ่านจำนวนหน้า PDF
        reader = PdfReader(pdf_file)
        total_pages = len(reader.pages)

        # Page range
        pr = self.page_range.get()
        if pr:
            try:
                start_page, end_page = map(int, pr.split("-"))
            except:
                messagebox.showerror("Error", "Page range format invalid (ex: 1-10)")
                return
        else:
            start_page, end_page = 1, total_pages
        if start_page < 1 or end_page > total_pages or start_page > end_page:
            messagebox.showerror("Error", "Page range is invalid")
            return

        # Export choice
        export_input = self.combo_export.get().split()[0]  # '1', '2', '3' หรือ '1,2,3'
        export_choice = [c.strip() for c in export_input.split(",")]

        # สร้างโฟลเดอร์สำหรับ image ถ้าจำเป็น
        if "1" in export_choice or "2" in export_choice or "3" in export_choice:
            images_folder = os.path.join(output_folder, f"{pdf_name}_images")
            os.makedirs(images_folder, exist_ok=True)
        else:
            images_folder = None

        # แปลง PDF → Images
        self.status.config(text="📖 Converting PDF to Images ...")
        self.root.update()
        image_files = []
        total_to_convert = end_page - start_page + 1
        for i in range(start_page, end_page + 1):
            pages = convert_from_path(pdf_file, dpi=150, first_page=i, last_page=i)
            if images_folder:
                image_path = os.path.join(images_folder, f"page_{i}.jpg")
                pages[0].save(image_path, "JPEG")
                image_files.append(image_path)
            percent = ((i - start_page + 1) / total_to_convert) * 100
            self.status.config(text=f"📖 Converting PDF: Page {i}/{end_page} ({percent:.1f}%)")
            self.root.update()

        # Export Word
        if "1" in export_choice:
            self.status.config(text="📝 Creating Word ...")
            self.root.update()
            word_file = os.path.join(output_folder, f"{pdf_name}.docx")
            doc = Document()
            section = doc.sections[0]
            section.page_width = Inches(8.27)
            section.page_height = Inches(11.69)
            section.left_margin = section.right_margin = 0
            section.top_margin = section.bottom_margin = 0
            for idx, img in enumerate(image_files, start=1):
                if idx > 1:
                    new_section = doc.add_section()
                    new_section.page_width = Inches(8.27)
                    new_section.page_height = Inches(11.69)
                    new_section.left_margin = new_section.right_margin = 0
                    new_section.top_margin = new_section.bottom_margin = 0
                    section = new_section
                paragraph = doc.add_paragraph()
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = paragraph.add_run()
                run.add_picture(img, width=section.page_width, height=section.page_height)
            doc.save(word_file)

        # Export PPT
        if "2" in export_choice:
            self.status.config(text="📊 Creating PowerPoint ...")
            self.root.update()
            pptx_file = os.path.join(output_folder, f"{pdf_name}.pptx")
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            for img in image_files:
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(pptx_file)

        # Export Images only
        if "3" in export_choice and images_folder:
            self.status.config(text=f"📂 Images saved at {images_folder}")
            self.root.update()

        messagebox.showinfo("Done", "🎉 Conversion Finished!")
        self.status.config(text="✅ Finished")

if __name__ == "__main__":
    root = tk.Tk()
    app = PDFConverterApp(root)
    root.mainloop()
